from __future__ import annotations

import hashlib
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any

import structlog
from docx import Document
from pypdf import PdfReader
from pptx import Presentation

from shared.models import EmbeddingRecord

log = structlog.get_logger()

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt", ".md", ".markdown"}


@dataclass(frozen=True)
class ExtractedDocument:
    path: Path
    root: Path
    extension: str
    text: str
    content_hash: str

    @property
    def corpus_group(self) -> str:
        return self.root.name or "legal_corpus"

    @property
    def relative_path(self) -> str:
        try:
            return str(self.path.relative_to(self.root))
        except ValueError:
            return self.path.name


def _normalize_text(text: str) -> str:
    cleaned = text.replace("\r\n", "\n").replace("\r", "\n")
    cleaned = re.sub(r"[ \t]+", " ", cleaned)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    return cleaned.strip()


def _extract_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    pages: list[str] = []
    for idx, page in enumerate(reader.pages, start=1):
        page_text = (page.extract_text() or "").strip()
        if page_text:
            pages.append(f"[page {idx}]\n{page_text}")
    return "\n\n".join(pages)


def _extract_docx(path: Path) -> str:
    doc = Document(str(path))
    blocks: list[str] = []
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            blocks.append(text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                blocks.append(" | ".join(cells))
    return "\n".join(blocks)


def _extract_pptx(path: Path) -> str:
    presentation = Presentation(str(path))
    blocks: list[str] = []
    for slide_idx, slide in enumerate(presentation.slides, start=1):
        slide_chunks: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = str(shape.text or "").strip()
                if text:
                    slide_chunks.append(text)
        if slide_chunks:
            blocks.append(f"[slide {slide_idx}]\n" + "\n".join(slide_chunks))
    return "\n\n".join(blocks)


def extract_document(path: Path, root: Path) -> ExtractedDocument | None:
    extension = path.suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        return None
    if extension in {".txt", ".md", ".markdown"}:
        raw = path.read_text(encoding="utf-8", errors="ignore")
    elif extension == ".pdf":
        raw = _extract_pdf(path)
    elif extension == ".docx":
        raw = _extract_docx(path)
    elif extension == ".pptx":
        raw = _extract_pptx(path)
    else:
        return None
    text = _normalize_text(raw)
    if not text:
        return None
    content_hash = hashlib.sha256(text.encode("utf-8")).hexdigest()
    return ExtractedDocument(
        path=path.resolve(),
        root=root.resolve(),
        extension=extension,
        text=text,
        content_hash=content_hash,
    )


def discover_documents(paths: list[str], max_files: int = 0) -> list[tuple[Path, Path]]:
    discovered: list[tuple[Path, Path]] = []
    for raw in paths:
        root = Path(raw).expanduser().resolve()
        if not root.exists():
            continue
        if root.is_file() and root.suffix.lower() in SUPPORTED_EXTENSIONS:
            discovered.append((root, root.parent))
            continue
        if not root.is_dir():
            continue
        for candidate in root.rglob("*"):
            if candidate.is_file() and candidate.suffix.lower() in SUPPORTED_EXTENSIONS:
                discovered.append((candidate.resolve(), root))
                if max_files > 0 and len(discovered) >= max_files:
                    return discovered
    discovered.sort(key=lambda item: str(item[0]))
    return discovered


def _chunk_text(text: str, max_chunk_chars: int = 1800, overlap_chars: int = 200) -> list[str]:
    paragraphs = [chunk.strip() for chunk in re.split(r"\n{2,}", text) if chunk.strip()]
    if not paragraphs:
        return []

    chunks: list[str] = []
    current = ""
    for paragraph in paragraphs:
        candidate = paragraph if not current else f"{current}\n\n{paragraph}"
        if len(candidate) <= max_chunk_chars:
            current = candidate
            continue
        if current:
            chunks.append(current)
        if len(paragraph) <= max_chunk_chars:
            current = paragraph
            continue
        start = 0
        while start < len(paragraph):
            end = min(start + max_chunk_chars, len(paragraph))
            chunks.append(paragraph[start:end].strip())
            if end >= len(paragraph):
                break
            start = max(0, end - overlap_chars)
        current = ""
    if current:
        chunks.append(current)
    return [chunk for chunk in chunks if chunk]


def build_embedding_records(
    *,
    paths: list[str],
    source: str = "legal_corpus",
    kind: str = "legal_chunk",
    max_chunk_chars: int = 1800,
    max_files: int = 0,
) -> tuple[list[EmbeddingRecord], dict[str, Any]]:
    records: list[EmbeddingRecord] = []
    discovered = discover_documents(paths, max_files=max_files)
    extracted_docs = 0
    skipped_docs = 0

    for file_path, root in discovered:
        try:
            extracted = extract_document(file_path, root)
        except Exception as exc:
            skipped_docs += 1
            log.warning("legal_corpus_extract_failed", path=str(file_path), error=str(exc))
            continue
        if extracted is None:
            skipped_docs += 1
            continue
        extracted_docs += 1
        chunks = _chunk_text(extracted.text, max_chunk_chars=max_chunk_chars)
        for chunk_idx, chunk in enumerate(chunks):
            source_id = hashlib.sha256(
                f"{extracted.relative_path}:{extracted.content_hash}:{chunk_idx}".encode("utf-8")
            ).hexdigest()[:24]
            records.append(
                EmbeddingRecord(
                    source=source,
                    kind=kind,
                    source_id=source_id,
                    content=chunk,
                    metadata={
                        "file_path": str(extracted.path),
                        "relative_path": extracted.relative_path,
                        "file_name": extracted.path.name,
                        "extension": extracted.extension,
                        "corpus_group": extracted.corpus_group,
                        "chunk_index": chunk_idx,
                        "content_hash": extracted.content_hash,
                    },
                )
            )

    stats = {
        "paths_requested": len(paths),
        "files_discovered": len(discovered),
        "files_extracted": extracted_docs,
        "files_skipped": skipped_docs,
        "chunks_built": len(records),
        "source": source,
        "kind": kind,
    }
    return records, stats
